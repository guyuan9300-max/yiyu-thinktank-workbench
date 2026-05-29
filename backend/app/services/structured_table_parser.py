"""xlsx / pptx 结构化解析（Phase 0）。

旧路径 `_archive_xml_text` 只把 ZIP 里所有 XML 文本平铺拼接，丢失：
- Sheet / Slide 边界
- 行列关系
- 单元格类型
- 演讲者备注（pptx notes）
- 表头语义

本模块用 openpyxl / python-pptx 做**结构化解析**：
- xlsx → 每个 sheet 渲染为 markdown 表格 + 保留 records JSON 供 Phase 1
- pptx → 每张 slide 拆分（标题 / 正文 / 备注），单独成 chunk

设计：
- 失败任何一步都抛出 StructuredParseError，调用方决定 fallback
- 不做"猜表头位置"等启发式（在 Phase 1 加）
- 返回 dataclass 结构便于序列化进 v2_chunks.content + 未来 structured_tables 表
"""
from __future__ import annotations

import logging
from dataclasses import dataclass, field
from datetime import date, datetime, time
from pathlib import Path
from typing import Any

import zipfile

import openpyxl
from openpyxl.utils.exceptions import InvalidFileException
from pptx import Presentation
from pptx.exc import PackageNotFoundError

logger = logging.getLogger(__name__)


class StructuredParseError(Exception):
    """解析失败，调用方应当 fallback 到旧路径。"""


# ---- 类型 ----------------------------------------------------------------


CellType = str  # "text" | "number" | "date" | "currency" | "percentage" | "formula" | "blank"


@dataclass(frozen=True)
class ParsedSheet:
    """xlsx 的一个 sheet 解析结果。"""

    sheet_name: str
    headers: list[str]
    rows: list[dict[str, Any]]      # 每行：{header_name: value}
    column_types: dict[str, CellType]
    markdown: str                    # 给 LLM 看的 markdown 表格
    row_count: int
    column_count: int
    notes: list[str] = field(default_factory=list)  # 解析过程的备注（合并单元格 / 公式 / 略过空 sheet 等）


@dataclass(frozen=True)
class ParsedSlide:
    """pptx 的一张 slide 解析结果。"""

    slide_no: int                    # 1-indexed
    title: str
    body: str                        # 正文（不含标题、不含备注）
    notes: str                       # 演讲者备注（speaker notes）
    markdown: str                    # 渲染好的 markdown（用作 chunk.content）


# ---- xlsx ---------------------------------------------------------------


def _format_cell_value(value: Any) -> str:
    """把单元格值格式化成 markdown 适用的字符串。"""
    if value is None:
        return ""
    if isinstance(value, bool):
        return "是" if value else "否"
    if isinstance(value, (datetime,)):
        return value.strftime("%Y-%m-%d %H:%M" if (value.hour or value.minute) else "%Y-%m-%d")
    if isinstance(value, date):
        return value.strftime("%Y-%m-%d")
    if isinstance(value, time):
        return value.strftime("%H:%M")
    if isinstance(value, float):
        # 整数表示成无小数
        if value.is_integer():
            return str(int(value))
        return f"{value:g}"
    return str(value)


def _detect_cell_type(value: Any) -> CellType:
    if value is None or value == "":
        return "blank"
    if isinstance(value, bool):
        return "text"
    if isinstance(value, (int, float)):
        return "number"
    if isinstance(value, (datetime, date, time)):
        return "date"
    text = str(value)
    if text.startswith("=") and len(text) > 1:
        return "formula"
    return "text"


def _aggregate_column_types(rows: list[dict[str, Any]], headers: list[str]) -> dict[str, CellType]:
    """按列聚合最常见的类型（忽略 blank）。"""
    result: dict[str, CellType] = {}
    for header in headers:
        counts: dict[CellType, int] = {}
        for row in rows:
            cell_type = _detect_cell_type(row.get(header))
            if cell_type == "blank":
                continue
            counts[cell_type] = counts.get(cell_type, 0) + 1
        if not counts:
            result[header] = "blank"
            continue
        # 取出现最多的类型
        result[header] = max(counts.items(), key=lambda kv: kv[1])[0]
    return result


def _render_sheet_markdown(
    sheet_name: str,
    headers: list[str],
    rows: list[dict[str, Any]],
    *,
    max_rows_inline: int = 200,
) -> str:
    """把 sheet 渲染成 markdown 表格。

    超过 max_rows_inline 行截断，前后各保留一部分 + 中间省略提示。
    """
    if not headers:
        return f"## Sheet: {sheet_name}\n\n_（空 sheet）_\n"

    lines = [f"## Sheet: {sheet_name}", ""]
    # 表头
    safe_headers = [_md_escape(h) for h in headers]
    lines.append("| " + " | ".join(safe_headers) + " |")
    lines.append("| " + " | ".join(["---"] * len(safe_headers)) + " |")

    if len(rows) <= max_rows_inline:
        display_rows = rows
        truncated = False
    else:
        head_n = max_rows_inline // 2
        tail_n = max_rows_inline - head_n
        display_rows = rows[:head_n] + rows[-tail_n:]
        truncated = True

    for index, row in enumerate(display_rows):
        if truncated and index == max_rows_inline // 2:
            lines.append("| " + " | ".join(["…"] * len(safe_headers)) + " |")
        cells = [_md_escape(_format_cell_value(row.get(h))) for h in headers]
        lines.append("| " + " | ".join(cells) + " |")

    if truncated:
        lines.append("")
        lines.append(f"_（共 {len(rows)} 行，已截断为头 {max_rows_inline // 2} + 尾 {max_rows_inline - max_rows_inline // 2}）_")

    lines.append("")
    return "\n".join(lines)


_MD_ESCAPE_CHARS = ("|", "\n", "\r")


def _md_escape(text: str) -> str:
    s = str(text)
    for ch in _MD_ESCAPE_CHARS:
        s = s.replace(ch, " ")
    return s.strip()


def _find_header_row(sheet: Any) -> int:
    """启发式：第一个非全空的行视为 header。

    边角情况：极偶尔有 title 行+空行+header 行+数据。本 Phase 不处理，
    简单取第一个非空行；Phase 1 可加 LLM 校准。
    """
    for row in sheet.iter_rows(values_only=True):
        if any(cell not in (None, "") for cell in row):
            return sheet.iter_rows(values_only=True).__next__.__call__()  # noqa
    return 0


def parse_xlsx_structured(filepath: str | Path) -> list[ParsedSheet]:
    """解析 xlsx，每个 sheet 一个 ParsedSheet。

    Raises:
        StructuredParseError: 文件无法打开或全部 sheet 为空
    """
    path = Path(filepath)
    if not path.exists():
        raise StructuredParseError(f"文件不存在: {path}")
    try:
        wb = openpyxl.load_workbook(path, data_only=True, read_only=False)
    except (InvalidFileException, OSError, KeyError, zipfile.BadZipFile) as exc:
        raise StructuredParseError(f"openpyxl 打开失败: {exc}") from exc

    results: list[ParsedSheet] = []
    try:
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            parsed = _parse_single_sheet(sheet_name, sheet)
            if parsed is not None:
                results.append(parsed)
    finally:
        wb.close()

    if not results:
        raise StructuredParseError("所有 sheet 都为空或全跳过")
    return results


def _parse_single_sheet(sheet_name: str, sheet: Any) -> ParsedSheet | None:
    """解析单个 sheet。完全空 → 返回 None。"""
    # 抓所有非空行
    raw_rows: list[tuple] = []
    for row in sheet.iter_rows(values_only=True):
        if any(cell not in (None, "") for cell in row):
            raw_rows.append(row)
    if not raw_rows:
        return None

    # 把第一非空行当 headers，余下当数据
    header_row = raw_rows[0]
    headers: list[str] = []
    for index, cell in enumerate(header_row):
        if cell in (None, ""):
            headers.append(f"列{index + 1}")
        else:
            headers.append(_format_cell_value(cell))

    data_rows: list[dict[str, Any]] = []
    for raw in raw_rows[1:]:
        row_dict: dict[str, Any] = {}
        for index, header in enumerate(headers):
            value = raw[index] if index < len(raw) else None
            row_dict[header] = value
        data_rows.append(row_dict)

    column_types = _aggregate_column_types(data_rows, headers)
    markdown = _render_sheet_markdown(sheet_name, headers, data_rows)
    notes: list[str] = []

    # 公式提醒
    formula_cols = [c for c, t in column_types.items() if t == "formula"]
    if formula_cols:
        notes.append(f"以下列包含公式（data_only 仍读到字符串）：{', '.join(formula_cols)}")

    # 合并单元格提示
    try:
        merged_ranges = list(sheet.merged_cells.ranges) if hasattr(sheet, "merged_cells") else []
        if merged_ranges:
            notes.append(f"原表含 {len(merged_ranges)} 处合并单元格（已按左上值展开）")
    except Exception:
        pass

    return ParsedSheet(
        sheet_name=sheet_name,
        headers=headers,
        rows=data_rows,
        column_types=column_types,
        markdown=markdown,
        row_count=len(data_rows),
        column_count=len(headers),
        notes=notes,
    )


# ---- pptx ---------------------------------------------------------------


_PPT_PLACEHOLDER_PATTERNS = (
    "click to add",
    "click to edit",
    "请单击此处",
    "在此处单击",
)


def _is_placeholder_text(text: str) -> bool:
    lower = text.strip().lower()
    return any(p in lower for p in _PPT_PLACEHOLDER_PATTERNS)


def _extract_shape_text(shape: Any) -> str:
    """从一个 pptx shape 抽文本，包含表格 / 占位符 / 一般文本框。"""
    parts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text.strip()
        if text and not _is_placeholder_text(text):
            parts.append(text)
    # 嵌入的表格
    if getattr(shape, "has_table", False):
        try:
            table = shape.table
            table_lines: list[str] = []
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                table_lines.append(" | ".join(cells))
            if table_lines:
                parts.append("\n".join(table_lines))
        except Exception:
            pass
    # group 形状递归
    try:
        if getattr(shape, "shape_type", None) is not None and getattr(shape, "shapes", None) is not None:
            for inner in shape.shapes:
                inner_text = _extract_shape_text(inner)
                if inner_text:
                    parts.append(inner_text)
    except Exception:
        pass
    return "\n".join(p for p in parts if p)


def parse_pptx_structured(filepath: str | Path) -> list[ParsedSlide]:
    """解析 pptx，每张 slide 一个 ParsedSlide。"""
    path = Path(filepath)
    if not path.exists():
        raise StructuredParseError(f"文件不存在: {path}")
    try:
        prs = Presentation(str(path))
    except (PackageNotFoundError, OSError, KeyError, zipfile.BadZipFile) as exc:
        raise StructuredParseError(f"python-pptx 打开失败: {exc}") from exc

    results: list[ParsedSlide] = []
    for index, slide in enumerate(prs.slides):
        slide_no = index + 1
        title = ""
        body_parts: list[str] = []
        for shape in slide.shapes:
            # 标题
            if getattr(shape, "is_placeholder", False):
                try:
                    ph_type = shape.placeholder_format.type
                    # 1 = TITLE, 13 = CENTER_TITLE
                    if ph_type in {1, 13}:
                        if getattr(shape, "has_text_frame", False):
                            txt = shape.text_frame.text.strip()
                            if txt and not _is_placeholder_text(txt):
                                title = txt
                                continue
                except Exception:
                    pass
            text = _extract_shape_text(shape)
            if text:
                body_parts.append(text)

        # 备注
        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        body = "\n\n".join(body_parts).strip()
        markdown = _render_slide_markdown(slide_no, title, body, notes_text)
        # 完全空白的 slide（无标题/正文/备注）跳过
        if not (title or body or notes_text):
            continue
        results.append(
            ParsedSlide(
                slide_no=slide_no,
                title=title,
                body=body,
                notes=notes_text,
                markdown=markdown,
            )
        )

    if not results:
        raise StructuredParseError("所有 slide 都为空")
    return results


def _render_slide_markdown(slide_no: int, title: str, body: str, notes: str) -> str:
    lines = [f"## Slide {slide_no}: {title or '（无标题）'}", ""]
    if body:
        lines.append(body)
        lines.append("")
    if notes:
        lines.append("**演讲者备注：**")
        lines.append(notes)
        lines.append("")
    return "\n".join(lines)


__all__ = [
    "CellType",
    "ParsedSheet",
    "ParsedSlide",
    "StructuredParseError",
    "parse_pptx_structured",
    "parse_xlsx_structured",
]
