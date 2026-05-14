"""VisualOcrRunner（Phase 1）。

接收 task payload 描述 "某文件的某个视觉单元"（如 pptx 第 7 张 slide 的所有嵌入图），
渲染/抽出图片 → 调 vision 模型 OCR → 把 markdown 写入 v2_chunks。

支持的 source_kind（Phase 1 只先实现 pptx_slide_images，其他留接口）：
- "pptx_slide_images": 从 .pptx ZIP 抽指定 slide 的嵌入图，每张 OCR 后合并
  payload 结构: {source_kind, source_path, source_v2_document_id,
                  source_client_id, region: {slide_no: int}, title}

将来扩展：
- "pdf_page": 从 pdf 渲染指定页为图 → OCR
- "image_file": 直接处理整张图片文件
- "xlsx_embedded": 从 xlsx 抽嵌入图（按 sheet+anchor 定位）

幂等性：runner 不删旧 v2_chunks。清理职责在 backfill 脚本（开始前一次性清掉）。
任务的 result_json 含 created_chunk_ids，重跑同一 task 会写多份 chunk，
依赖外层 input_hash 去重避免重复 enqueue。
"""
from __future__ import annotations

import base64
import json
import logging
import time
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any
from uuid import uuid4

from app.db import Database
from app.services.knowledge_v2 import normalize_text, tokenize

logger = logging.getLogger(__name__)


SOURCE_PPTX_SLIDE = "pptx_slide_images"


def _now_iso() -> str:
    return datetime.now().isoformat(timespec="seconds")


def _new_chunk_id() -> str:
    return f"v2chunk_visual_{uuid4().hex[:10]}"


def _parse_payload(task: dict[str, object]) -> dict[str, Any]:
    """task.payload_json 是入队时塞的执行参数。"""
    raw = task.get("payload_json") or "{}"
    try:
        payload = json.loads(raw) if isinstance(raw, str) else dict(raw)
    except json.JSONDecodeError:
        payload = {}
    return payload if isinstance(payload, dict) else {}


def _extract_pptx_slide_images(
    pptx_path: Path,
    slide_no: int,
) -> tuple[str, str, str, list[bytes]]:
    """从 .pptx 抽指定 slide 的：title / body / notes / 所有嵌入图 PNG/JPEG bytes。

    返回：(title, body_text, notes_text, [image_bytes, ...])

    实现细节：
    - python-pptx 给我们 slide 的 shapes，包含 picture shape 和 text frame
    - picture shape 的 image bytes 通过 shape.image.blob 取
    - 同时拿到 slide 自身的文字（title placeholder / 其他 text frame / notes）
    """
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)
    if slide_no < 1 or slide_no > len(slides):
        raise ValueError(f"slide_no={slide_no} 越界（pptx 共 {len(slides)} 页）")
    slide = slides[slide_no - 1]

    title = ""
    body_parts: list[str] = []
    images: list[bytes] = []

    for shape in slide.shapes:
        # 标题占位符
        if getattr(shape, "is_placeholder", False):
            try:
                ph_type = shape.placeholder_format.type
                # 1=TITLE, 13=CENTER_TITLE
                if ph_type in {1, 13} and getattr(shape, "has_text_frame", False):
                    txt = shape.text_frame.text.strip()
                    if txt and not title:
                        title = txt
                        continue
            except Exception:
                pass
        # 一般文字
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                body_parts.append(text)
        # 嵌入表格
        if getattr(shape, "has_table", False):
            try:
                rows: list[str] = []
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    body_parts.append("\n".join(rows))
            except Exception:
                pass
        # 图片
        if getattr(shape, "shape_type", None) == 13:  # PICTURE
            try:
                blob = shape.image.blob
                if blob and len(blob) > 1024:  # 跳过 <1KB 的小装饰
                    images.append(blob)
            except Exception:
                pass

    body = "\n\n".join(p for p in body_parts if p).strip()

    # 演讲者备注
    notes = ""
    try:
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        pass

    return title, body, notes, images


def _detect_image_mime(blob: bytes) -> str:
    """简易 magic-byte 嗅探，覆盖 PPT 常见图片格式。"""
    if blob.startswith(b"\x89PNG"):
        return "image/png"
    if blob.startswith(b"\xff\xd8\xff"):
        return "image/jpeg"
    if blob.startswith(b"GIF8"):
        return "image/gif"
    if blob.startswith(b"RIFF") and blob[8:12] == b"WEBP":
        return "image/webp"
    return "image/png"  # 默认


def _ocr_images_with_ai(
    ai_service: Any,
    *,
    title: str,
    slide_no: int,
    images: list[bytes],
) -> str:
    """对一个 slide 的所有嵌入图依次 OCR，合并 markdown。"""
    if not images or not hasattr(ai_service, "generate_visual_markdown"):
        return ""
    parts: list[str] = []
    for index, blob in enumerate(images, start=1):
        b64 = base64.b64encode(blob).decode()
        mime = _detect_image_mime(blob)
        try:
            md = ai_service.generate_visual_markdown(
                title=title or "PPT slide",
                page_number=slide_no,
                image_base64=b64,
                mime_type=mime,
                source_kind="PPT 图片",
            )
        except Exception as exc:  # noqa: BLE001
            logger.warning(
                "visual_ocr slide=%d image=%d failed: %s",
                slide_no, index, exc,
            )
            md = ""
        md = (md or "").strip()
        if md and len(md) >= 8:  # 太短当噪音丢
            tag = f"### 图 {index}" if len(images) > 1 else ""
            parts.append((tag + "\n" + md).strip())
    return "\n\n".join(parts).strip()


def _compose_chunk_markdown(
    *,
    slide_no: int,
    title: str,
    body: str,
    notes: str,
    visual_md: str,
) -> str:
    """把 slide 的所有信息（标题/正文/备注/图 OCR）拼成一份 markdown，写入 chunk.content。"""
    sections: list[str] = [f"## Slide {slide_no}: {title or '（无标题）'}"]
    if body:
        sections.append(body)
    if visual_md:
        sections.append("**图片识别：**\n" + visual_md)
    if notes:
        sections.append("**演讲者备注：**\n" + notes)
    return "\n\n".join(sections).strip()


def _write_section_and_chunk(
    db: Database,
    *,
    v2_document_id: str,
    section_index: int,
    chunk_index: int,
    content: str,
    section_label: str,
) -> str:
    """每 slide 建一个 v2_section，下挂一条 v2_chunk。返回 chunk_id。"""
    now = _now_iso()
    normalized = normalize_text(content)
    searchable = " ".join(tokenize(normalized))[:4000]

    section_id = f"v2section_visual_{uuid4().hex[:10]}"
    db.execute(
        """
        INSERT INTO v2_sections (
            id, v2_document_id, section_index, title, content,
            searchable_text, char_count, created_at
        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        """,
        (
            section_id, v2_document_id, section_index, section_label,
            content, searchable, len(content), now,
        ),
    )

    chunk_id = _new_chunk_id()
    db.execute(
        """
        INSERT INTO v2_chunks (
            id, v2_document_id, v2_section_id, chunk_index, section_label,
            content, searchable_text, char_count, created_at
        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
        """,
        (
            chunk_id, v2_document_id, section_id, chunk_index, section_label,
            content, searchable, len(content), now,
        ),
    )
    return chunk_id


def _next_section_index(db: Database, v2_document_id: str) -> int:
    row = db.fetchone(
        "SELECT MAX(section_index) AS m FROM v2_sections WHERE v2_document_id = ?",
        (v2_document_id,),
    )
    if not row or row["m"] is None:
        return 0
    return int(row["m"]) + 1


def _next_chunk_index(db: Database, v2_document_id: str) -> int:
    row = db.fetchone(
        "SELECT MAX(chunk_index) AS m FROM v2_chunks WHERE v2_document_id = ?",
        (v2_document_id,),
    )
    if not row or row["m"] is None:
        return 0
    return int(row["m"]) + 1


# ──────────────────────────────────────────────────────────────────
# Runner 主入口
# ──────────────────────────────────────────────────────────────────


def process(db: Database, ai_service: Any, task: dict[str, object]) -> dict[str, object]:
    """供 local_model_optimizer.run_due_local_model_tasks dispatch 调用。

    成功返回 result_json dict（会被写入 local_model_tasks.result_json）。
    失败抛异常，外层负责 mark_failed + 重试计数。
    """
    payload = _parse_payload(task)
    source_kind = str(payload.get("source_kind") or "")
    source_path = str(payload.get("source_path") or "")
    v2_document_id = str(payload.get("source_v2_document_id") or task.get("knowledge_document_id") or "")
    title = str(payload.get("title") or "")

    if source_kind != SOURCE_PPTX_SLIDE:
        raise RuntimeError(f"visual_ocr_runner 暂不支持 source_kind={source_kind!r}")
    if not source_path or not Path(source_path).exists():
        raise FileNotFoundError(f"pptx 不存在: {source_path}")
    if not v2_document_id:
        raise RuntimeError("payload 缺 source_v2_document_id")

    region = payload.get("region") or {}
    if not isinstance(region, dict):
        region = {}
    slide_no = int(region.get("slide_no") or 0)
    if slide_no <= 0:
        raise RuntimeError("payload.region.slide_no 缺失或非法")

    started = time.time()
    pptx_title, body, notes, images = _extract_pptx_slide_images(Path(source_path), slide_no)

    visual_md = _ocr_images_with_ai(
        ai_service,
        title=title or pptx_title or Path(source_path).stem,
        slide_no=slide_no,
        images=images,
    )

    content = _compose_chunk_markdown(
        slide_no=slide_no,
        title=pptx_title,
        body=body,
        notes=notes,
        visual_md=visual_md,
    )

    chunk_id = None
    if content and len(content.strip()) >= 16:
        chunk_id = _write_section_and_chunk(
            db,
            v2_document_id=v2_document_id,
            section_index=_next_section_index(db, v2_document_id),
            chunk_index=_next_chunk_index(db, v2_document_id),
            content=content,
            section_label=f"Slide {slide_no}",
        )

    duration_ms = int((time.time() - started) * 1000)
    return {
        "source_kind": source_kind,
        "source_path": source_path,
        "region": {"slide_no": slide_no},
        "title": pptx_title or title,
        "image_count": len(images),
        "visual_md_chars": len(visual_md),
        "total_content_chars": len(content),
        "chunk_id": chunk_id,
        "duration_ms": duration_ms,
        "ran_at": _now_iso(),
    }


__all__ = ["process", "SOURCE_PPTX_SLIDE"]
