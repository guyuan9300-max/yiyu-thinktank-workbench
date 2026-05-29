"""xlsx / pptx 结构化解析测试（Phase 0）。"""
from __future__ import annotations

import sys
from datetime import date, datetime
from pathlib import Path

import openpyxl
import pytest
from pptx import Presentation
from pptx.util import Inches

sys.path.append(str(Path(__file__).resolve().parents[1]))

from app.services.structured_table_parser import (
    ParsedSheet,
    ParsedSlide,
    StructuredParseError,
    parse_pptx_structured,
    parse_xlsx_structured,
)


# ---- xlsx ---------------------------------------------------------------


def _build_xlsx(path: Path, sheet_data: dict[str, list[list]]) -> Path:
    wb = openpyxl.Workbook()
    # 删默认空 sheet
    default = wb.active
    wb.remove(default)
    for sheet_name, rows in sheet_data.items():
        ws = wb.create_sheet(title=sheet_name)
        for row in rows:
            ws.append(row)
    wb.save(path)
    return path


@pytest.mark.unit
def test_parse_xlsx_basic_table(tmp_path: Path) -> None:
    p = tmp_path / "budget.xlsx"
    _build_xlsx(p, {
        "Q1预算": [
            ["项目类别", "预算金额", "已花费", "剩余"],
            ["调研经费", 50000, 32000, 18000],
            ["受益人活动", 80000, 45000, 35000],
            ["培训费", 30000, 12000, 18000],
        ],
    })
    sheets = parse_xlsx_structured(p)
    assert len(sheets) == 1
    sheet = sheets[0]
    assert sheet.sheet_name == "Q1预算"
    assert sheet.headers == ["项目类别", "预算金额", "已花费", "剩余"]
    assert sheet.row_count == 3
    assert sheet.column_count == 4
    assert sheet.column_types["预算金额"] == "number"
    assert sheet.column_types["项目类别"] == "text"
    assert sheet.rows[0]["项目类别"] == "调研经费"
    assert sheet.rows[0]["预算金额"] == 50000
    # markdown 里能看到表头 + 数据
    assert "## Sheet: Q1预算" in sheet.markdown
    assert "| 项目类别 | 预算金额 | 已花费 | 剩余 |" in sheet.markdown
    assert "50000" in sheet.markdown
    assert "调研经费" in sheet.markdown


@pytest.mark.unit
def test_parse_xlsx_multi_sheet(tmp_path: Path) -> None:
    p = tmp_path / "multi.xlsx"
    _build_xlsx(p, {
        "立项": [["项目", "负责人"], ["A", "张三"], ["B", "李四"]],
        "预算": [["项目", "金额"], ["A", 50000], ["B", 30000]],
        "成果": [["项目", "受益人数"], ["A", 120], ["B", 80]],
    })
    sheets = parse_xlsx_structured(p)
    assert len(sheets) == 3
    names = {s.sheet_name for s in sheets}
    assert names == {"立项", "预算", "成果"}
    budget = next(s for s in sheets if s.sheet_name == "预算")
    assert budget.column_types["金额"] == "number"


@pytest.mark.unit
def test_parse_xlsx_skips_empty_sheets(tmp_path: Path) -> None:
    p = tmp_path / "with_empty.xlsx"
    _build_xlsx(p, {
        "数据": [["A", "B"], [1, 2]],
        "空": [],
    })
    sheets = parse_xlsx_structured(p)
    assert len(sheets) == 1
    assert sheets[0].sheet_name == "数据"


@pytest.mark.unit
def test_parse_xlsx_dates_are_formatted(tmp_path: Path) -> None:
    p = tmp_path / "dated.xlsx"
    _build_xlsx(p, {
        "时间": [
            ["事件", "日期"],
            ["上线", date(2026, 6, 1)],
            ["复盘", datetime(2026, 8, 15, 14, 30)],
        ],
    })
    sheets = parse_xlsx_structured(p)
    sheet = sheets[0]
    assert sheet.column_types["日期"] == "date"
    # markdown 里日期格式应该是 YYYY-MM-DD
    assert "2026-06-01" in sheet.markdown
    assert "2026-08-15" in sheet.markdown


@pytest.mark.unit
def test_parse_xlsx_handles_missing_header_cells(tmp_path: Path) -> None:
    """表头有空 cell → 用 列N 占位，不应崩。"""
    p = tmp_path / "sparse.xlsx"
    _build_xlsx(p, {
        "稀疏": [
            ["A", None, "C"],
            [1, 2, 3],
        ],
    })
    sheets = parse_xlsx_structured(p)
    assert sheets[0].headers == ["A", "列2", "C"]


@pytest.mark.unit
def test_parse_xlsx_truncates_huge_table(tmp_path: Path) -> None:
    p = tmp_path / "huge.xlsx"
    rows = [["序号", "值"]]
    for i in range(1, 401):
        rows.append([i, i * 10])
    _build_xlsx(p, {"big": rows})
    sheets = parse_xlsx_structured(p)
    sheet = sheets[0]
    assert sheet.row_count == 400
    # markdown 应当被截断
    assert "截断" in sheet.markdown


@pytest.mark.unit
def test_parse_xlsx_missing_file_raises() -> None:
    with pytest.raises(StructuredParseError):
        parse_xlsx_structured("/tmp/__no_such_file__.xlsx")


@pytest.mark.unit
def test_parse_xlsx_corrupted_file_raises(tmp_path: Path) -> None:
    bad = tmp_path / "broken.xlsx"
    bad.write_bytes(b"this is not an xlsx")
    with pytest.raises(StructuredParseError):
        parse_xlsx_structured(bad)


# ---- pptx ---------------------------------------------------------------


def _build_pptx(path: Path, slides: list[dict]) -> Path:
    """slides: list of {'title': str, 'body': str, 'notes': str}"""
    prs = Presentation()
    blank_layout = prs.slide_layouts[5]  # title only
    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)
        # 标题
        if slide_data.get("title"):
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = slide_data["title"]
        # 正文（在标题下方加一个文本框）
        if slide_data.get("body"):
            left = top = Inches(1)
            width = Inches(8)
            height = Inches(4)
            box = slide.shapes.add_textbox(left, top, width, height)
            box.text_frame.text = slide_data["body"]
        # 备注
        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]
    prs.save(path)
    return path


@pytest.mark.unit
def test_parse_pptx_basic_slides(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    _build_pptx(p, [
        {"title": "项目背景", "body": "客户在 2025 年启动战略升级。"},
        {"title": "核心数据", "body": "服务人次 1200，预算 50 万。", "notes": "这个数据是估算的，正式汇报别说"},
    ])
    slides = parse_pptx_structured(p)
    assert len(slides) == 2
    assert slides[0].slide_no == 1
    assert slides[0].title == "项目背景"
    assert "客户在 2025 年启动战略升级。" in slides[0].body
    assert slides[1].title == "核心数据"
    # 备注里的关键信息保留
    assert "这个数据是估算的" in slides[1].notes
    # markdown 里也能看到备注
    assert "演讲者备注" in slides[1].markdown


@pytest.mark.unit
def test_parse_pptx_skips_blank_slides(tmp_path: Path) -> None:
    p = tmp_path / "with_blank.pptx"
    _build_pptx(p, [
        {"title": "有内容", "body": "正文"},
        {"title": "", "body": "", "notes": ""},  # 全空 → 跳过
    ])
    slides = parse_pptx_structured(p)
    assert len(slides) == 1
    assert slides[0].title == "有内容"


@pytest.mark.unit
def test_parse_pptx_corrupted_raises(tmp_path: Path) -> None:
    bad = tmp_path / "broken.pptx"
    bad.write_bytes(b"not a pptx")
    with pytest.raises(StructuredParseError):
        parse_pptx_structured(bad)


@pytest.mark.unit
def test_parse_pptx_missing_file_raises() -> None:
    with pytest.raises(StructuredParseError):
        parse_pptx_structured("/tmp/__no_such_file__.pptx")


# ---- 集成：与 knowledge_v2 ingest 链对接 -----------------------------------


@pytest.mark.integration
def test_extract_document_uses_structured_xlsx(tmp_path: Path) -> None:
    """ingest 链对 xlsx 走 structured parser，每个 sheet 成为一个 section。"""
    from app.services.knowledge_v2 import extract_document_with_metadata

    p = tmp_path / "client_dashboard.xlsx"
    _build_xlsx(p, {
        "项目概览": [["项目", "状态"], ["A", "进行中"]],
        "预算": [["科目", "金额"], ["调研", 50000]],
    })
    doc = extract_document_with_metadata(p)
    assert doc.metadata.parse_status == "ready"
    section_titles = {s["title"] for s in doc.sections}
    assert "Sheet · 项目概览" in section_titles
    assert "Sheet · 预算" in section_titles
    # markdown 表格应当在 cleaned text 里
    assert "项目" in doc.text
    assert "50000" in doc.text


@pytest.mark.integration
def test_extract_document_uses_structured_pptx(tmp_path: Path) -> None:
    """ingest 链对 pptx 走 structured parser，每张 slide 一个 section + 备注保留。"""
    from app.services.knowledge_v2 import extract_document_with_metadata

    p = tmp_path / "deck.pptx"
    _build_pptx(p, [
        {"title": "战略 Q1", "body": "重点：组织升级", "notes": "汇报时强调风险"},
    ])
    doc = extract_document_with_metadata(p)
    assert doc.metadata.parse_status == "ready"
    assert any("Slide 1" in s["title"] for s in doc.sections)
    assert "汇报时强调风险" in doc.text  # 备注被收录


@pytest.mark.integration
def test_extract_document_falls_back_to_legacy_when_corrupted(tmp_path: Path) -> None:
    """损坏的 xlsx 走 fallback：返回 failed 但不抛错。"""
    from app.services.knowledge_v2 import extract_document_with_metadata

    bad = tmp_path / "broken.xlsx"
    bad.write_bytes(b"not a real xlsx")
    doc = extract_document_with_metadata(bad)
    assert doc.metadata.parse_status == "failed"
